import os
from pptx import Presentation
import fitz  # PyMuPDF

def extract_slide_text(file_path):
    if file_path.endswith(".pptx"):
        prs = Presentation(file_path)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    cleaned = shape.text.strip()
                    if cleaned:
                        all_text.append(cleaned)
        return "\n".join(all_text)
    
    elif file_path.endswith(".pdf"):
        doc = fitz.open(file_path)
        all_text = []
        for page in doc:
            text = page.get_text()
            if text.strip():
                all_text.append(text.strip())
        return "\n\n".join(all_text)
    
    else:
        raise ValueError("Unsupported file type. Use .pdf or .pptx")    # ✅ Correct

if __name__ == "__main__":
    file_path = input("Enter the path to your .pptx or .pdf file: ")
    try:
        extracted_text = extract_slide_text(file_path)
        print("\n--- Extracted Text ---\n")
        print(extracted_text)
    except Exception as e:
        print(f"Error: {e}")

